import re

from docx import Document
from pptx import Presentation

# ── helpers ───────────────────────────────────────────────────────────────────

_TEMPLATE_RE = re.compile(r'^%\w+%$')
_DIGITS_RE   = re.compile(r'^\d+$')
_SKIP_TEXT   = {'Submit', 'Next', 'Question'}

def _clean(text):
    """Remove Storyline line-number prefixes and collapse segment splits."""
    text = re.sub(r'\n\n\d', ' ', text)
    text = re.sub(r'^\d', '', text)
    return text.strip()


def _is_item_state(type_str):
    """True if this row type represents an answer-state variant (not Normal)."""
    tl = type_str.lower()
    if 'alttext' in tl or 'alt text' in tl:
        return False  # image alt-text rows, not answer choices
    return ('selected state' in tl or 'sw state'      in tl or
            'crt state'      in tl or 'crct state'    in tl or
            'inc state'      in tl or 'ict state'     in tl or
            'drop correct'   in tl or 'drop incorrect' in tl)


def _is_skip(text):
    """True if text is a UI label, template variable, or pure digit — not content."""
    return (not text or text in _SKIP_TEXT or
            bool(_TEMPLATE_RE.match(text)) or bool(_DIGITS_RE.match(text)))


def _set_cell(cell, text, ref_run=None):
    for para in cell.paragraphs:
        for run in para.runs:
            run.text = ''
    if not text:
        return
    para = cell.paragraphs[0]
    if para.runs:
        para.runs[0].text = text
    else:
        run = para.add_run(text)
        if ref_run:
            if ref_run.font.name:
                run.font.name = ref_run.font.name
            if ref_run.font.size:
                run.font.size = ref_run.font.size
            if ref_run.font.bold is not None:
                run.font.bold = ref_run.font.bold
            if ref_run.font.italic is not None:
                run.font.italic = ref_run.font.italic


def _split_proportional(new_text, orig_full, orig_seg1):
    """Split new_text at a word boundary proportional to orig_seg1/orig_full."""
    if not orig_full or not new_text:
        return new_text, ''
    ratio = len(orig_seg1) / max(len(orig_full), 1)
    pos   = int(len(new_text) * ratio)
    left  = new_text.rfind(' ', 0, pos)
    right = new_text.find(' ', pos)
    if left == -1 and right == -1:
        return new_text, ''
    sp = left if left != -1 and (right == -1 or (pos - left) <= (right - pos)) else right
    return new_text[:sp].strip(), new_text[sp:].strip()


# ── question type detection ───────────────────────────────────────────────────

def _detect_type(data_rows):
    # Lowercase + exclude AltText rows for state detection
    types = [r[1].lower() for r in data_rows
             if len(r) > 1 and 'alttext' not in r[1].lower() and 'alt text' not in r[1].lower()]
    has_drop      = any('drop correct' in t or 'drop incorrect' in t for t in types)
    # any state marker that labels answer choices (case-insensitive)
    has_opt_state = any(
        'crt state'      in t or 'crct state'     in t or
        'inc state'      in t or 'ict state'      in t or
        'selected state' in t or 'sw state'       in t
        for t in types
    )
    if has_drop:
        return 'DND'
    if has_opt_state:
        item_texts = {_clean(r[2]) for r in data_rows
                      if _is_item_state(r[1]) and len(r) > 2}
        return 'TF' if len(item_texts) <= 2 else 'MCQ'
    return None


def _stem_and_options(data_rows, q_type=None):
    # Build the set of option texts from any answer-state row
    item_texts = {_clean(r[2]) for r in data_rows
                  if _is_item_state(r[1]) and len(r) > 2}

    options_seen, options_set = [], set()
    stem = None

    for row in data_rows:
        if len(row) < 3 or row[1] == 'Slide name':
            continue
        tl = row[1].lower()
        if 'normal state' not in tl:
            continue
        if 'alttext' in tl or 'alt text' in tl:
            continue  # image alt-text, not a translatable choice
        cleaned = _clean(row[2])
        if _is_skip(cleaned):
            continue
        if cleaned in item_texts:
            if cleaned not in options_set:
                options_seen.append(cleaned)
                options_set.add(cleaned)
        elif stem is None:
            stem = cleaned

    return stem, options_seen


# ── PPTX parser ───────────────────────────────────────────────────────────────

def get_pptx_questions(pptx_path):
    """
    Returns list of dicts:
      { id, stem, options, preview }
    Only slides with a question stem + at least one option are included.
    """
    prs = Presentation(pptx_path)
    result = []

    for idx, slide in enumerate(prs.slides):
        texts = [
            shape.text_frame.text.strip()
            for shape in slide.shapes
            if shape.has_text_frame and shape.text_frame.text.strip()
        ]
        if not texts:
            continue

        block = max(texts, key=len)
        # \x0b = PowerPoint soft-return; treat it as a line break
        lines = [ln.strip() for ln in block.replace('\x0b', '\n').split('\n') if ln.strip()]

        if len(lines) < 2:
            continue

        stem = lines[0]
        # strip A. / a. / A) / a) / 1. / 1) option-letter prefixes
        options = [re.sub(r'^[A-Da-d1-9][.)]\s*', '', ln).strip() for ln in lines[1:]]

        result.append({
            'id':       idx + 1,
            'stem':     stem,
            'options':  options,
            'preview':  stem[:120] + ('…' if len(stem) > 120 else ''),
        })

    return result


# ── Word parser — format detection ───────────────────────────────────────────

def _is_single_table_doc(tables):
    """True if the doc uses the new flat SingleTable export (no per-slide metadata tables)."""
    for t in tables:
        if t.rows and 'Slide ID' in t.rows[0].cells[0].text:
            return False
    return True


def _find_content_table(tables):
    for t in tables:
        if t.rows and 'ID \U0001f512' in t.rows[0].cells[0].text:
            return t
    return None


def _parse_single_table(content_table):
    """
    Parse a SingleTable-format DOCX into question dicts.
    Returns list of { slide_id, slide_name, type, stem, options, preview }.
    """
    all_rows = list(content_table.rows)
    slides   = []     # list of (slide_name, data_rows_text)
    cur_slide = None
    cur_data  = []

    for ri, row in enumerate(all_rows):
        if ri == 0:
            continue  # skip header
        cells = [c.text.strip() for c in row.cells]
        if len(cells) < 2:
            continue
        rtype = cells[1]
        src   = cells[2] if len(cells) > 2 else ''

        if 'Slide name' in rtype:
            if cur_slide is not None:
                slides.append((cur_slide, cur_data))
            cur_slide = src
            cur_data  = []
        elif 'Scene name' in rtype:
            if cur_slide is not None:
                slides.append((cur_slide, cur_data))
            cur_slide = None
            cur_data  = []
        elif cur_slide is not None:
            cur_data.append(cells)

    if cur_slide is not None:
        slides.append((cur_slide, cur_data))

    result = []
    for slide_name, data_rows in slides:
        q_type = _detect_type(data_rows)
        if not q_type:
            continue
        stem, opts = _stem_and_options(data_rows, q_type)
        if not opts and not stem:
            continue
        result.append({
            'slide_id':   slide_name,
            'slide_name': slide_name,
            'type':       q_type,
            'stem':       stem or '',
            'options':    opts,
            'preview':    (stem or '')[:120] + ('…' if len(stem or '') > 120 else ''),
        })

    return result


def _parse_single_table_for_apply(content_table):
    """
    Like _parse_single_table but also returns (row_start, row_end) indices
    into the content_table for write-back in apply_mappings.
    Returns list of { slide_id, slide_name, type, stem, options, data_rows,
                       table, row_start, row_end }.
    """
    all_rows = list(content_table.rows)
    slides   = []   # (slide_name, data_rows_text, row_start, row_end)
    cur_slide = None
    cur_start = None
    cur_data  = []

    for ri, row in enumerate(all_rows):
        if ri == 0:
            continue
        cells = [c.text.strip() for c in row.cells]
        if len(cells) < 2:
            continue
        rtype = cells[1]
        src   = cells[2] if len(cells) > 2 else ''

        if 'Slide name' in rtype:
            if cur_slide is not None:
                slides.append((cur_slide, cur_data, cur_start, ri))
            cur_slide = src
            cur_start = ri + 1
            cur_data  = []
        elif 'Scene name' in rtype:
            if cur_slide is not None:
                slides.append((cur_slide, cur_data, cur_start, ri))
            cur_slide = None
            cur_data  = []
        elif cur_slide is not None:
            cur_data.append(cells)

    if cur_slide is not None:
        slides.append((cur_slide, cur_data, cur_start, len(all_rows)))

    result = []
    for slide_name, data_rows, row_start, row_end in slides:
        q_type = _detect_type(data_rows)
        if not q_type:
            continue
        stem, opts = _stem_and_options(data_rows, q_type)
        result.append({
            'slide_id':   slide_name,
            'slide_name': slide_name,
            'type':       q_type,
            'stem':       stem or '',
            'options':    opts,
            'data_rows':  data_rows,
            'table':      content_table,
            'row_start':  row_start,
            'row_end':    row_end,
        })

    return result


# ── Word parser — public API ──────────────────────────────────────────────────

def get_word_questions(docx_path):
    """
    Returns list of dicts:
      { slide_id, slide_name, type, stem, options, preview }
    Handles both paired-table and SingleTable export formats.
    Only question slides (MCQ / TF / DND) are included.
    """
    doc    = Document(docx_path)
    tables = doc.tables

    if _is_single_table_doc(tables):
        ct = _find_content_table(tables)
        return _parse_single_table(ct) if ct else []

    # ── paired-table format (original) ────────────────────────────────────────
    result = []
    i = 0

    while i < len(tables):
        row0 = [c.text.strip() for c in tables[i].rows[0].cells] if tables[i].rows else []

        if len(row0) == 2 and 'Slide ID' in row0[0]:
            slide_id = row0[1]

            if i + 1 < len(tables):
                ct_rows = [[c.text.strip() for c in r.cells] for r in tables[i + 1].rows]

                if ct_rows and 'ID \U0001f512' in ct_rows[0][0]:
                    data = ct_rows[1:]
                    slide_name = next(
                        (r[2] for r in data if len(r) > 2 and r[1] == 'Slide name'), ''
                    )
                    q_type = _detect_type(data)

                    if q_type:
                        stem, opts = _stem_and_options(data, q_type)
                        result.append({
                            'slide_id':   slide_id,
                            'slide_name': slide_name,
                            'type':       q_type,
                            'stem':       stem or '',
                            'options':    opts,
                            'preview':    (stem or '')[:120] + ('…' if len(stem or '') > 120 else ''),
                        })

                    i += 2
                    continue

        i += 1

    return result


# ── apply user-defined mappings ───────────────────────────────────────────────

def apply_mappings(docx_path, pptx_path, mappings):
    """
    mappings: list of { word_slide_id: str, pptx_id: int }
    Opens docx, applies each mapping by writing PPTX content into Translation column,
    returns (Document, results_list).
    results_list entries: { slide_name, filled, mismatch }
    """
    pptx_list  = get_pptx_questions(pptx_path)
    pptx_by_id = {q['id']: q for q in pptx_list}

    doc    = Document(docx_path)
    tables = doc.tables
    word_qs = []

    if _is_single_table_doc(tables):
        ct = _find_content_table(tables)
        if ct:
            word_qs = _parse_single_table_for_apply(ct)
    else:
        # ── paired-table format ───────────────────────────────────────────────
        i = 0
        while i < len(tables):
            row0 = [c.text.strip() for c in tables[i].rows[0].cells] if tables[i].rows else []

            if len(row0) == 2 and 'Slide ID' in row0[0]:
                slide_id = row0[1]

                if i + 1 < len(tables):
                    ct      = tables[i + 1]
                    ct_rows = [[c.text.strip() for c in r.cells] for r in ct.rows]

                    if ct_rows and 'ID \U0001f512' in ct_rows[0][0]:
                        data = ct_rows[1:]
                        q_type = _detect_type(data)

                        if q_type:
                            stem, opts = _stem_and_options(data, q_type)
                            slide_name = next(
                                (r[2] for r in data if len(r) > 2 and r[1] == 'Slide name'), ''
                            )
                            word_qs.append({
                                'slide_id':   slide_id,
                                'slide_name': slide_name,
                                'type':       q_type,
                                'stem':       stem or '',
                                'options':    opts,
                                'data_rows':  data,
                                'table':      ct,
                                'row_start':  1,
                            })

                    i += 2
                    continue

            i += 1

    word_by_id = {q['slide_id']: q for q in word_qs}

    results = []
    for m in mappings:
        wq     = word_by_id.get(m['word_slide_id'])
        pptx_q = pptx_by_id.get(int(m['pptx_id']))

        if not wq or not pptx_q:
            continue

        filled, mismatch = _fill(wq, pptx_q)
        results.append({'slide_name': wq['slide_name'], 'filled': filled, 'mismatch': mismatch})

    return doc, results


def _fill(word_q, pptx_entry):
    """
    Write pptx_entry stem+options into Translation column of word_q['table'].
    Always writes — no similarity check (user explicitly chose this mapping).
    Returns (filled_count, option_mismatch_count).
    """
    text_map = {}
    if word_q['stem'] and pptx_entry['stem']:
        text_map[word_q['stem']] = pptx_entry['stem']
    elif pptx_entry['stem']:
        # Template slide: %question% is a runtime placeholder; write PPTX stem there
        text_map['%question%'] = pptx_entry['stem']
    for w_opt, p_opt in zip(word_q['options'], pptx_entry['options']):
        text_map[w_opt] = p_opt
    mismatch = abs(len(word_q['options']) - len(pptx_entry['options']))

    table     = word_q['table']
    data      = word_q['data_rows']
    row_start = word_q.get('row_start', 1)
    row_end   = word_q.get('row_end', None)

    table_rows = list(table.rows)[row_start:row_end]
    filled     = 0
    seen_key   = {}

    for ri, table_row in enumerate(table_rows):
        if ri >= len(data):
            break

        row = data[ri]
        if len(row) < 4:
            continue
        row_id, row_type, src, cur_trans = row[0], row[1], row[2], row[3]

        if row_type == 'Slide name' or not src:
            continue

        src_cell = table_row.cells[2]
        ref_run  = (src_cell.paragraphs[0].runs[0]
                    if src_cell.paragraphs and src_cell.paragraphs[0].runs
                    else None)

        is_multi = '\n\n' in src
        key      = (row_id, row_type)

        if is_multi:
            if key not in seen_key:
                cleaned  = _clean(src)
                new_text = text_map.get(cleaned)
                if new_text:
                    seg1, seg2 = _split_proportional(new_text, cleaned, cur_trans)
                    _set_cell(table_row.cells[3], seg1, ref_run)
                    seen_key[key] = seg2
                    filled += 1
                else:
                    seen_key[key] = None
            else:
                seg2 = seen_key.pop(key)
                if seg2 is not None:
                    _set_cell(table_row.cells[3], seg2, ref_run)
                    filled += 1
        else:
            cleaned  = _clean(src)
            new_text = text_map.get(cleaned)
            if new_text:
                _set_cell(table_row.cells[3], new_text, ref_run)
                filled += 1

    return filled, mismatch
